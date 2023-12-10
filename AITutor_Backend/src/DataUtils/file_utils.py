from pptx import Presentation
import glob
import os
import datetime

def read_pptx_file(pptx_file):
    content = ''
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content+=shape.text+'\n'
    return content

def save_training_data(output_dir, input_data, output_data):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{output_dir}output_{timestamp}.txt"
    try:
        with open(filename, 'w') as file:
            file.write(input_data + "\n[->LLM_OUTPUT]\n" + output_data)
    except IOError as e:
        print(f"Error writing training data to file: {e}")